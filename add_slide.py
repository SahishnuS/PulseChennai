from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation('PulseChennai.pptx')

# Add a blank slide (layout 6 is usually blank)
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# Add title
txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "Recent Architecture Implementation (AI & Cloud)"
p.font.size = Pt(36)
p.font.bold = True

# Add content
contentBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
tf_content = contentBox.text_frame
tf_content.word_wrap = True

p = tf_content.paragraphs[0]
p.text = "Here are the critical engineering updates completed today:"
p.font.size = Pt(24)
p.font.bold = True

points = [
    "1. AI Inference Pipeline Integration: Built the /api/predict route and mounted the PyTorch InferencePipeline to serve real-time Spatial GNN predictions to the dashboard.",
    "2. PyTorch Graph Neural Network (GNN) Training: Resolved dynamic PyG tracing bugs and successfully trained the SpatialGNN on historical data, exporting pulse_gnn_best.pt.",
    "3. Cloud Parquet Data Lake (Batch Layer): Implemented the local Data Lake adapter and generated 25,000 synthetic historical trajectories for offline model training.",
    "4. HMM Map-Matching Activation: Expanded the Viterbi road network to 40 major segments and successfully wired probabilistic map-snapping into the live Redis message handler.",
    "5. Collaborative Telemetry Simulator: Built a real-time passenger simulator to continuously push ground-truth GPS pings, successfully activating the Ghost Bus fusion fail-safe."
]

for pt in points:
    p = tf_content.add_paragraph()
    p.text = pt
    p.font.size = Pt(18)
    p.level = 1

prs.save('PulseChennai_Final.pptx')
print("Saved PulseChennai_Final.pptx")
