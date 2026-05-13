import os
from pptx import Presentation

prs = Presentation('PulseChennai.pptx')

replacements = {
    # Slide 4
    "Predictive Dead Reckoning": "AI Graph Neural Networks (GNN)",
    "Speed + geometry prediction": "PyTorch GAT + LSTM deep learning models",
    "To stop \"Ghost Buses,\" our engine uses AI to detect if a bus is Stationary vs. Active. If a bus hasn't moved for 5 minutes at a terminal, it is automatically hidden from the \"Nearby\" view.": "To recover \"Ghost Buses,\" our AI uses PyTorch GNNs trained on Cloud Data Lakes to predictively project the bus forward using learned traffic patterns.",
    "5-minute inactivity threshold": "AI-powered predictive routing",
    
    # Slide 5
    "Processing Layer": "AI & Cloud Processing Layer",
    "Spatial Intelligence Engine": "Cloud-Native AI Inference Pipeline",
    "Uber H3 Library": "Redis Speed Layer & Uber H3",
    "Spatial binning & hexagonal indexing": "Real-time in-memory feature store & spatial mesh",
    "Kalman Filters": "PyTorch Graph Neural Networks",
    "Noise reduction & motion prediction": "Batch-trained GAT models on Parquet Data Lakes",
}

for slide in prs.slides:
    for shape in slide.shapes:
        if hasattr(shape, "text_frame"):
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    for old_text, new_text in replacements.items():
                        if old_text in run.text:
                            run.text = run.text.replace(old_text, new_text)
                            
                # Fallback in case the text is split across multiple runs
                text = shape.text_frame.text
                for old_text, new_text in replacements.items():
                    if old_text in text:
                        # Clear text frame and replace
                        # We only do this if it wasn't caught by the run replacement
                        # but it destroys formatting, so we try to be careful.
                        pass

# Second pass for full paragraph replacements to handle text split across runs
for slide in prs.slides:
    for shape in slide.shapes:
        if hasattr(shape, "text_frame"):
            for old_text, new_text in replacements.items():
                if old_text in shape.text_frame.text:
                    # If we find it in the text frame but it wasn't replaced in runs
                    if old_text in shape.text_frame.text:
                        # Find the paragraph
                        for p in shape.text_frame.paragraphs:
                            if old_text in p.text:
                                p.text = p.text.replace(old_text, new_text)

prs.save('PulseChennai_Updated.pptx')
print("Saved PulseChennai_Updated.pptx")
