from pptx import Presentation

prs = Presentation('PulseChennai.pptx')

with open("ppt_content.txt", "w", encoding="utf-8") as f:
    for i, slide in enumerate(prs.slides):
        f.write(f"--- Slide {i+1} ---\n")
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                f.write(f"[{shape.name}]: {shape.text}\n")
